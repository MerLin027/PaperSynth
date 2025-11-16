import os
import json
import fitz  # PyMuPDF for PDF processing
from fastapi import FastAPI, UploadFile, HTTPException
from fastapi import Request
from fastapi import Header, Depends
from fastapi.middleware.cors import CORSMiddleware
import asyncio
from fastapi.staticfiles import StaticFiles
from dotenv import load_dotenv
# Conditional imports will be done later when needed
StableDiffusionXLPipeline = None
try:
    from elevenlabs import ElevenLabs
except ImportError:
    ElevenLabs = None
import requests
from pptx import Presentation
import uvicorn
from huggingface_hub import login
import google.generativeai as genai
import logging 
from fpdf import FPDF
import torch
from datetime import datetime
from typing import Optional, Union
import uuid
import time
import shutil
from starlette.middleware.base import BaseHTTPMiddleware
from starlette.responses import Response
from fastapi.responses import FileResponse
import hmac
import hashlib
from urllib.parse import urlencode

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[logging.FileHandler("app.log"), logging.StreamHandler()]
)

# Load environment variables from backend.env
load_dotenv('backend.env')

# Create a directory for temporary files
TEMP_DIR = "temp_files"
os.makedirs(TEMP_DIR, exist_ok=True)

# Cleanup policy (configurable)
TEMP_TTL_HOURS = int(os.getenv("TEMP_TTL_HOURS", "24"))
TEMP_SIZE_CAP_GB = float(os.getenv("TEMP_SIZE_CAP_GB", "1"))
TTL_SECONDS = max(1, TEMP_TTL_HOURS) * 60 * 60
SIZE_CAP_BYTES = max(0.1, TEMP_SIZE_CAP_GB) * (1024 ** 3)
MAX_UPLOAD_BYTES = 10 * 1024 * 1024  # 10 MB
ALLOWED_SUMMARY_LENGTHS = {"short", "medium", "long"}
MAX_PDF_PAGES = int(os.getenv("MAX_PDF_PAGES", "100"))
MAX_TEXT_CHARS = int(os.getenv("MAX_TEXT_CHARS", "800000"))

# Get API keys securely
ELEVENLABS_API_KEY = os.getenv("ELEVENLABS_API_KEY")
HUGGINGFACE_API_TOKEN = os.getenv("HUGGINGFACE_API_TOKEN")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
API_AUTH_TOKEN = os.getenv("API_AUTH_TOKEN")
RATE_LIMIT_PER_MINUTE = int(os.getenv("RATE_LIMIT_PER_MINUTE", "10"))
CONCURRENCY_LIMIT = int(os.getenv("CONCURRENCY_LIMIT", "2"))
ALLOWED_CORS_ORIGINS = os.getenv("ALLOWED_CORS_ORIGINS", "http://localhost:3000").split(",")
ENABLE_SDXL = os.getenv("ENABLE_SDXL", "true").lower() == "true"
ENABLE_TTS = os.getenv("ENABLE_TTS", "true").lower() == "true"
SIGNED_DOWNLOADS = os.getenv("SIGNED_DOWNLOADS", "false").lower() == "true"
DOWNLOAD_SIGNING_KEY = os.getenv("DOWNLOAD_SIGNING_KEY", "")

# Initialize Gemini client
if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)
    # Using gemini-flash-latest (always uses the newest stable flash model)
    # This avoids quota issues with experimental models
    gemini_model = genai.GenerativeModel('gemini-flash-latest')
else:
    gemini_model = None

# Ensure required API keys are set
missing_keys = []
if not GEMINI_API_KEY:
    missing_keys.append("GEMINI_API_KEY")
if ENABLE_TTS and not ELEVENLABS_API_KEY:
    missing_keys.append("ELEVENLABS_API_KEY")

# Memory management constants
MIN_MEMORY_GB = float(os.getenv("MIN_MEMORY_GB", "1.0"))
MEMORY_WARNING_THRESHOLD = float(os.getenv("MEMORY_WARNING_THRESHOLD", "0.85"))

if missing_keys:
    logging.error(f"Required API keys are missing: {', '.join(missing_keys)}")
    logging.warning("Application will run with limited functionality")
else:
    logging.info("All required API keys are configured")

# Environment Validation System
def validate_environment():
    """Comprehensive environment validation at startup"""
    validation_results = {
        "gemini": {"status": "unknown", "details": ""},
        "elevenlabs": {"status": "unknown", "details": ""},
        "memory": {"status": "unknown", "details": ""}
    }
    
    # 1. Gemini API Check
    try:
        if not GEMINI_API_KEY:
            validation_results["gemini"] = {"status": "error", "details": "GEMINI_API_KEY not found in environment variables"}
        elif not gemini_model:
            validation_results["gemini"] = {"status": "error", "details": "Gemini client initialization failed"}
        else:
            # Test API connection with a simple request
            test_response = gemini_model.generate_content("Test connection")
            if test_response and test_response.text:
                validation_results["gemini"] = {"status": "ok", "details": "Gemini API accessible and responsive"}
            else:
                validation_results["gemini"] = {"status": "warning", "details": "Gemini API connected but no response received"}
    except Exception as e:
        validation_results["gemini"] = {"status": "error", "details": f"Gemini API not accessible: {str(e)}"}
    
    # 2. ElevenLabs API Key Validation
    try:
        if not ElevenLabs:
            validation_results["elevenlabs"] = {"status": "error", "details": "ElevenLabs package not available. Install with: pip install elevenlabs"}
        elif ELEVENLABS_API_KEY and ELEVENLABS_API_KEY != "your-elevenlabs-api-key-here":
            client = ElevenLabs(api_key=ELEVENLABS_API_KEY)
            # Try a minimal API call to validate key
            voices = client.voices.get_all()
            validation_results["elevenlabs"] = {"status": "ok", "details": f"ElevenLabs API key valid, {len(voices.voices)} voices available"}
        else:
            validation_results["elevenlabs"] = {"status": "warning", "details": "ElevenLabs API key not configured (using placeholder)"}
    except Exception as e:
        validation_results["elevenlabs"] = {"status": "error", "details": f"ElevenLabs API validation failed: {str(e)}"}
    
    # 3. Hugging Face Token Validation
    try:
        if HUGGINGFACE_API_TOKEN and HUGGINGFACE_API_TOKEN != "your-huggingface-token-here":
            # Actually authenticate with Hugging Face if SDXL is enabled
            if ENABLE_SDXL:
                try:
                    login(HUGGINGFACE_API_TOKEN)
                    validation_results["huggingface"] = {"status": "ok", "details": "Hugging Face token authenticated successfully"}
                except Exception as e:
                    validation_results["huggingface"] = {"status": "error", "details": f"Hugging Face login failed: {str(e)}"}
            else:
                validation_results["huggingface"] = {"status": "ok", "details": "Hugging Face token present (SDXL disabled, not validated)"}
        else:
            if ENABLE_SDXL:
                validation_results["huggingface"] = {"status": "warning", "details": "Hugging Face token not configured but SDXL enabled"}
            else:
                validation_results["huggingface"] = {"status": "ok", "details": "Hugging Face token not needed (SDXL disabled)"}
    except Exception as e:
        validation_results["huggingface"] = {"status": "error", "details": f"Hugging Face authentication failed: {str(e)}"}
    
    # 4. System Resources Check
    try:
        import psutil
        memory = psutil.virtual_memory()
        disk = psutil.disk_usage(TEMP_DIR)
        
        system_info = {
            "memory_total_gb": round(memory.total / (1024**3), 1),
            "memory_available_gb": round(memory.available / (1024**3), 1),
            "disk_free_gb": round(disk.free / (1024**3), 1),
            "cuda_available": torch.cuda.is_available(),
            "cuda_devices": torch.cuda.device_count() if torch.cuda.is_available() else 0
        }
        
        # Check minimum requirements
        warnings = []
        if system_info["memory_available_gb"] < 4:
            warnings.append("Low memory (<4GB available)")
        if system_info["disk_free_gb"] < 2:
            warnings.append("Low disk space (<2GB free)")
        if not system_info["cuda_available"]:
            warnings.append("No CUDA GPU detected (will use CPU)")
            
        status = "warning" if warnings else "ok"
        details = f"Memory: {system_info['memory_available_gb']}/{system_info['memory_total_gb']}GB, Disk: {system_info['disk_free_gb']}GB free, CUDA: {system_info['cuda_available']}"
        if warnings:
            details += f" | Warnings: {', '.join(warnings)}"
            
        validation_results["system"] = {"status": status, "details": details}
    except ImportError:
        validation_results["system"] = {"status": "warning", "details": "psutil not available, cannot check system resources"}
    except Exception as e:
        validation_results["system"] = {"status": "error", "details": f"System check failed: {str(e)}"}
    
    # 5. Model Loading Test (lightweight check)
    try:
        # Test if we can at least import the models without loading them
        try:
            from diffusers import StableDiffusionXLPipeline
            validation_results["models"] = {"status": "ok", "details": "Model imports successful, lazy loading will be used"}
        except ImportError as e:
            validation_results["models"] = {"status": "warning", "details": f"Diffusers not available: {str(e)}. Visual generation will be disabled."}
    except Exception as e:
        validation_results["models"] = {"status": "error", "details": f"Model import failed: {str(e)}"}
    
    # Summarize results
    errors = []
    warnings = []
    success = True
    
    for component, result in validation_results.items():
        if result["status"] == "error":
            errors.append(component)
            success = False
        elif result["status"] == "warning":
            warnings.append(component)
    
    validation_results["success"] = success
    validation_results["errors"] = errors
    validation_results["warnings"] = warnings
    
    return validation_results

def log_validation_results(results):
    """Log validation results in a readable format"""
    logging.info("=== Environment Validation Results ===")
    
    for component, result in results.items():
        if component in ['success', 'errors', 'warnings']:
            continue
            
        status_symbol = "PASS" if result['success'] else "WARN" if result.get('warning') else "FAIL"
        logging.info(f"{status_symbol} {component.upper()}: {result['details']}")
    
    errors = results.get('errors', [])
    warnings = results.get('warnings', [])
    
    if errors:
        logging.warning(f"Critical failures detected in: {', '.join(errors)}")
        logging.warning("Application may not function properly. Please check configuration.")
    
    if warnings:
        logging.info(f"Non-critical issues detected in: {', '.join(warnings)}")
        logging.info("These may cause degraded functionality but won't prevent startup.")

# Environment validation system is ready
# Validation will run automatically on startup via the @app.on_event("startup") handler

app = FastAPI()

# Serve generated files as static assets
app.mount("/static", StaticFiles(directory=TEMP_DIR), name="static")

# CORS for browser clients (limit to known origins)
app.add_middleware(
    CORSMiddleware,
    allow_origins=[o.strip() for o in ALLOWED_CORS_ORIGINS if o.strip()],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Singleton cache for Stable Diffusion pipeline with lazy loading
_SDXL_PIPE = None
_SDXL_LOADING = False
_SDXL_LOAD_LOCK = asyncio.Lock()

# Memory management utilities
def get_memory_usage():
    """Get current memory usage statistics"""
    try:
        import psutil
        process = psutil.Process()
        memory_info = process.memory_info()
        return {
            "rss_mb": round(memory_info.rss / (1024 * 1024), 1),
            "vms_mb": round(memory_info.vms / (1024 * 1024), 1),
            "percent": round(process.memory_percent(), 2)
        }
    except:
        return {"rss_mb": 0, "vms_mb": 0, "percent": 0}

def cleanup_gpu_memory():
    """Clean up GPU memory when possible"""
    try:
        if torch.cuda.is_available():
            torch.cuda.empty_cache()
            logging.info("GPU memory cache cleared")
    except Exception as e:
        logging.warning(f"Failed to clear GPU memory: {e}")

async def lazy_load_stable_diffusion():
    """Lazy load Stable Diffusion with proper async handling"""
    global _SDXL_PIPE, _SDXL_LOADING
    
    if _SDXL_PIPE is not None:
        return _SDXL_PIPE
    
    async with _SDXL_LOAD_LOCK:
        # Double-check pattern
        if _SDXL_PIPE is not None:
            return _SDXL_PIPE
            
        if _SDXL_LOADING:
            # Wait for another thread to finish loading
            while _SDXL_LOADING:
                await asyncio.sleep(0.1)
            return _SDXL_PIPE
        
        _SDXL_LOADING = True
        try:
            logging.info("Starting lazy load of Stable Diffusion XL...")
            memory_before = get_memory_usage()
            
            # Run the actual loading in a thread pool to avoid blocking
            def _load_sdxl():
                return initialize_stable_diffusion()
            
            # Use thread pool for CPU-intensive loading
            loop = asyncio.get_event_loop()
            _SDXL_PIPE = await loop.run_in_executor(None, _load_sdxl)
            
            memory_after = get_memory_usage()
            logging.info(f"SDXL loaded successfully. Memory usage: {memory_before['rss_mb']}MB -> {memory_after['rss_mb']}MB")
            return _SDXL_PIPE
            
        except Exception as e:
            logging.error(f"Failed to lazy load SDXL: {e}")
            raise
        finally:
            _SDXL_LOADING = False

# Global concurrency semaphore
_CONC_SEM = asyncio.Semaphore(CONCURRENCY_LIMIT)

# In-memory rate limit buckets
_RATE_BUCKETS = {}

def _get_dir_size_bytes(path: str) -> int:
    total = 0
    for root, _dirs, files in os.walk(path):
        for name in files:
            fp = os.path.join(root, name)
            try:
                total += os.path.getsize(fp)
            except OSError:
                pass
    return total

def _cleanup_temp_dir():
    try:
        # Collect subdirectories with their modified times and sizes
        subdirs = []
        now = time.time()
        total_size = 0
        for name in os.listdir(TEMP_DIR):
            full = os.path.join(TEMP_DIR, name)
            if os.path.isdir(full):
                try:
                    mtime = os.path.getmtime(full)
                except OSError:
                    mtime = now
                size = _get_dir_size_bytes(full)
                total_size += size
                subdirs.append({"path": full, "mtime": mtime, "size": size})
        
        # TTL deletion
        for entry in subdirs:
            if now - entry["mtime"] > TTL_SECONDS:
                try:
                    shutil.rmtree(entry["path"], ignore_errors=True)
                except Exception:
                    pass

        # Recompute after TTL deletion
        total_size = 0
        kept = []
        for name in os.listdir(TEMP_DIR):
            full = os.path.join(TEMP_DIR, name)
            if os.path.isdir(full):
                size = _get_dir_size_bytes(full)
                total_size += size
                try:
                    mtime = os.path.getmtime(full)
                except OSError:
                    mtime = now
                kept.append({"path": full, "mtime": mtime, "size": size})

        # Size cap deletion: delete oldest until under cap
        if total_size > SIZE_CAP_BYTES:
            kept.sort(key=lambda e: e["mtime"])  # oldest first
            for entry in kept:
                if total_size <= SIZE_CAP_BYTES:
                    break
                try:
                    shutil.rmtree(entry["path"], ignore_errors=True)
                    total_size -= entry["size"]
                except Exception:
                    pass
    except Exception as e:
        logging.warning(f"Cleanup task error: {e}")

def _start_cleanup_background_task():
    import threading
    def _runner():
        while True:
            _cleanup_temp_dir()
            time.sleep(60 * 60)  # hourly
    t = threading.Thread(target=_runner, daemon=True)
    t.start()

@app.on_event("startup")
def _on_startup():
    # Run environment validation on startup
    validation_result = validate_environment()
    if not validation_result["success"]:
        logging.error(f"Environment validation failed: {validation_result['errors']}")
        logging.warning("Application may not function properly with these issues")
    else:
        logging.info("Environment validation passed successfully")
    
    _start_cleanup_background_task()

class RequestIdLoggingMiddleware(BaseHTTPMiddleware):
    async def dispatch(self, request, call_next):
        rid = uuid.uuid4().hex
        request.state.request_id = rid
        # Redact Authorization from logs
        auth = request.headers.get("authorization")
        if auth:
            redacted_headers = dict(request.headers)
            redacted_headers["authorization"] = "[REDACTED]"
            logging.info(f"[{rid}] {request.method} {request.url.path} headers={ {k: redacted_headers[k] for k in ['host','authorization'] if k in redacted_headers} }")
        else:
            logging.info(f"[{rid}] {request.method} {request.url.path}")
        response: Response = await call_next(request)
        response.headers["X-Request-ID"] = rid
        return response

app.add_middleware(RequestIdLoggingMiddleware)

def initialize_stable_diffusion():
    """Initialize Stable Diffusion with optimal settings"""
    try:
        global _SDXL_PIPE, StableDiffusionXLPipeline
        if _SDXL_PIPE is not None:
            return _SDXL_PIPE

        # Try to import StableDiffusionXLPipeline if not already done
        if StableDiffusionXLPipeline is None:
            try:
                from diffusers import StableDiffusionXLPipeline
            except ImportError as e:
                raise HTTPException(
                    status_code=500,
                    detail=f"StableDiffusionXLPipeline not available. Error: {str(e)}. Please fix dependency issues."
                )

        if torch.cuda.is_available():
            # GPU path with memory-efficient settings
            pipe = StableDiffusionXLPipeline.from_pretrained(
                "stabilityai/stable-diffusion-xl-base-1.0",
                torch_dtype=torch.float16,
                use_safetensors=True,
                variant="fp16"
            )
            pipe = pipe.to("cuda")

            if torch.cuda.get_device_properties(0).total_memory >= 8 * (1024 ** 3):  # 8GB or more
                try:
                    pipe.enable_xformers_memory_efficient_attention()
                except Exception as e:
                    logging.warning(f"Could not enable xformers: {e}")
                    pipe.enable_attention_slicing()
            else:
                pipe.enable_attention_slicing()
                pipe.enable_sequential_cpu_offload()
            _SDXL_PIPE = pipe
            return _SDXL_PIPE
        else:
            # CPU fallback
            if StableDiffusionXLPipeline is None:
                try:
                    from diffusers import StableDiffusionXLPipeline
                except ImportError as e:
                    raise HTTPException(
                        status_code=500,
                        detail=f"StableDiffusionXLPipeline not available. Error: {str(e)}. Please fix dependency issues."
                    )
            logging.warning("CUDA not available. Falling back to CPU for Stable Diffusion; generation will be slower.")
            pipe = StableDiffusionXLPipeline.from_pretrained(
                "stabilityai/stable-diffusion-xl-base-1.0",
                torch_dtype=torch.float32,
                use_safetensors=True
            )
            pipe = pipe.to("cpu")
            pipe.enable_attention_slicing()
            _SDXL_PIPE = pipe
            return _SDXL_PIPE
    except Exception as e:
        logging.error(f"Error initializing Stable Diffusion: {str(e)}")
        # Provide more specific error messages based on common issues
        if "CUDA" in str(e):
            raise HTTPException(
                status_code=500,
                detail="GPU initialization failed. Please ensure CUDA is properly installed and a compatible GPU is available."
            )
        elif "memory" in str(e).lower():
            raise HTTPException(
                status_code=500,
                detail="Insufficient GPU memory. Try reducing batch size or image dimensions."
            )
        else:
            raise HTTPException(status_code=500, detail=f"Failed to initialize Stable Diffusion: {str(e)}")

@app.get("/")
def read_root():
    return {"message": "FastAPI server is running!"}
# Health endpoint
@app.get("/health")
def health():
    try:
        # Run full environment validation
        validation_result = validate_environment()
        
        # Get memory usage
        memory_usage = get_memory_usage()
        
        return {
            "status": "healthy" if validation_result["success"] else "degraded",
            "temp_dir": TEMP_DIR,
            "rate_limit_per_minute": RATE_LIMIT_PER_MINUTE,
            "concurrency_limit": CONCURRENCY_LIMIT,
            "features": {
                "sdxl_enabled": ENABLE_SDXL,
                "tts_enabled": ENABLE_TTS,
                "signed_downloads": SIGNED_DOWNLOADS
            },
            "validation": validation_result,
            "memory": memory_usage
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Health check failed: {str(e)}")

# Status endpoint
@app.get("/status/{request_id}")
def status(request_id: str, request: Request):
    base_url = str(request.base_url).rstrip('/')
    req_dir = os.path.join(TEMP_DIR, request_id)
    if not os.path.isdir(req_dir):
        raise HTTPException(status_code=404, detail="request_id not found")
    def exists(name: str) -> bool:
        return os.path.exists(os.path.join(req_dir, name))
    def build_url(name: str):
        if SIGNED_DOWNLOADS and DOWNLOAD_SIGNING_KEY:
            expires = int(time.time()) + 10 * 60
            msg = f"{request_id}:{name}:{expires}".encode()
            sig = hmac.new(DOWNLOAD_SIGNING_KEY.encode(), msg, hashlib.sha256).hexdigest()
            qs = urlencode({"rid": request_id, "file": name, "exp": expires, "sig": sig})
            return f"{base_url}/download?{qs}"
        return f"{base_url}/static/{request_id}/{name}"
    return {
        "request_id": request_id,
        "summary_pdf": build_url("summary.pdf") if exists("summary.pdf") else None,
        "graphical_abstract": build_url("graphical_abstract.png") if exists("graphical_abstract.png") else None,
        "voiceover": build_url("voiceover.mp3") if exists("voiceover.mp3") else None,
        "presentation": build_url("presentation.pptx") if exists("presentation.pptx") else None,
    }

# Signed download endpoint
@app.get("/download")
def download_signed(rid: str, file: str, exp: int, sig: str):
    if not (SIGNED_DOWNLOADS and DOWNLOAD_SIGNING_KEY):
        raise HTTPException(status_code=404, detail="Downloads not signed")
    now = int(time.time())
    if now > exp:
        raise HTTPException(status_code=410, detail="Link expired")
    # Validate signature
    msg = f"{rid}:{file}:{exp}".encode()
    expected = hmac.new(DOWNLOAD_SIGNING_KEY.encode(), msg, hashlib.sha256).hexdigest()
    if not hmac.compare_digest(expected, sig):
        raise HTTPException(status_code=403, detail="Invalid signature")
    # Prevent path traversal
    if "/" in file or ".." in file:
        raise HTTPException(status_code=400, detail="Invalid file parameter")
    path = os.path.join(TEMP_DIR, rid, file)
    if not os.path.exists(path):
        raise HTTPException(status_code=404, detail="File not found")
    return FileResponse(path)

def require_bearer_token(authorization: Optional[str] = Header(default=None)):
    # If no token configured, allow all (development convenience)
    if not API_AUTH_TOKEN:
        return True
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing or invalid Authorization header")
    token = authorization.removeprefix("Bearer ").strip()
    if token != API_AUTH_TOKEN:
        raise HTTPException(status_code=403, detail="Invalid token")
    return True

def _client_key(request: Request, authorization: Optional[str]) -> str:
    # Prefer per-token limiting if available, else fall back to client IP
    if authorization and authorization.startswith("Bearer "):
        return f"tok:{authorization.removeprefix('Bearer ').strip()}"
    # Use direct client IP. Behind proxies, configure a proxy to pass real IP.
    client_ip = request.client.host if request.client else "unknown"
    return f"ip:{client_ip}"

def _rate_limit_check(key: str, now: float) -> bool:
    # Token bucket refill per second
    capacity = RATE_LIMIT_PER_MINUTE
    refill_rate_per_sec = RATE_LIMIT_PER_MINUTE / 60.0
    bucket = _RATE_BUCKETS.get(key)
    if not bucket:
        _RATE_BUCKETS[key] = {"tokens": capacity - 1, "last": now}
        return True
    elapsed = max(0.0, now - bucket["last"])
    bucket["last"] = now
    bucket["tokens"] = min(capacity, bucket["tokens"] + elapsed * refill_rate_per_sec)
    if bucket["tokens"] >= 1.0:
        bucket["tokens"] -= 1.0
        return True
    return False

def enforce_rate_limit(request: Request, authorization: Optional[str] = Header(default=None)):
    if RATE_LIMIT_PER_MINUTE <= 0:
        return True
    key = _client_key(request, authorization)
    if _rate_limit_check(key, time.time()):
        return True
    raise HTTPException(status_code=429, detail="Too many requests. Please slow down.")

# Extract text from PDF
def extract_text_from_pdf(pdf_path):
    try:
        with fitz.open(pdf_path) as doc:
            texts = []
            total_pages = len(doc)
            pages = min(total_pages, MAX_PDF_PAGES)
            for i in range(pages):
                texts.append(doc[i].get_text("text"))
                if sum(len(t) for t in texts) >= MAX_TEXT_CHARS:
                    break
            combined = "".join(texts)
            if len(combined) > MAX_TEXT_CHARS:
                combined = combined[:MAX_TEXT_CHARS]
            return combined, total_pages
    except Exception as e:
        # Fail fast with a clear message for corrupted/invalid PDFs
        msg = str(e).lower()
        if "cannot open" in msg or "broken" in msg or "corrupt" in msg or "invalid" in msg:
            logging.error(f"Invalid or corrupted PDF: {e}")
            raise HTTPException(status_code=400, detail="Invalid or corrupted PDF file")
        logging.error(f"PDF Extraction Error: {str(e)}")
        raise HTTPException(status_code=400, detail="Failed to read PDF file")

# Structured error helper
def _error_response(code: str, message: str, request_id: Optional[str], hint: Optional[str] = None):
    payload = {
        "code": code,
        "message": message,
    }
    if request_id:
        payload["request_id"] = request_id
    if hint:
        payload["hint"] = hint
    return payload

# Primary Gemini Summary Function
def gemini_summary(text, summary_length="medium"):
    try:
        if not gemini_model:
            raise HTTPException(status_code=500, detail="Gemini client not initialized - check GEMINI_API_KEY")
        
        # Create appropriate prompt for different lengths
        length_instructions = {
            "short": "Provide a concise summary in 200-300 words",
            "medium": "Provide a detailed summary in 400-600 words", 
            "long": "Provide a comprehensive summary in 800-1000 words"
        }
        length_instruction = length_instructions.get(summary_length, "Provide a detailed summary in 400-600 words")
        
        prompt = f"""You are an expert research paper analyst. Analyze this paper and provide a structured, engaging summary. {length_instruction}

Research Paper Content:
{text}

Create a comprehensive summary including:

**Overview:** Provide a 2-3 sentence summary of the research question, methods, and main findings.

**Key Findings:** List the main discoveries, results, and notable data. Cite specific claims with page/section references where relevant.

**Methodology:** Describe the research approach, methods, and experimental design used.

**Conclusions:** Summarize the main conclusions drawn by the authors and any limitations acknowledged.

**Implications:** Discuss the broader impact, real-world applications, significance for the field, and future research directions.

Use clear language while maintaining technical accuracy. Include notable quotes that capture essential insights. Format with clear section headers and bullet points where appropriate to make it easy to understand and compelling to read."""

        response = gemini_model.generate_content(prompt)
        
        if response and response.text:
            logging.info("Gemini summary generated successfully.")
            return response.text
        else:
            raise Exception("No response generated from Gemini")
            
    except Exception as e:
        error_msg = str(e)
        logging.error(f"Gemini Summary Error: {error_msg}")
        
        # Provide more specific error messages
        if "API_KEY_INVALID" in error_msg or "API Key not found" in error_msg:
            raise HTTPException(status_code=500, detail="Invalid or missing Gemini API key. Please check your GEMINI_API_KEY in environment variables.")
        elif "quota" in error_msg.lower() or "limit" in error_msg.lower():
            raise HTTPException(status_code=429, detail="Gemini API quota exceeded. Please check your billing and usage limits.")
        elif "429" in error_msg or "rate" in error_msg.lower():
            raise HTTPException(status_code=429, detail="Gemini API rate limit exceeded. Please try again later.")
        else:
            raise HTTPException(status_code=500, detail=f"Gemini API Error: {error_msg}")

def format_summary_sections(summary):
    """Format the summary into structured sections using simple text parsing"""
    sections = {
        "Key Findings": [],
        "Methodology": [],
        "Conclusions": [],
        "Implications": []
    }
    
    try:
        # Split the summary into lines and clean them
        lines = [line.strip() for line in summary.split('\n') if line.strip()]
        
        current_section = None
        
        # Parse the summary looking for section headers and content
        for line in lines:
            line_lower = line.lower()
            
            # Check for section headers (including markdown format)
            if any(keyword in line_lower for keyword in ['**key finding', 'key finding', '**finding', 'finding']):
                current_section = "Key Findings"
                continue
            elif any(keyword in line_lower for keyword in ['**methodolog', 'methodolog', '**method', 'method', 'approach']):
                current_section = "Methodology"
                continue
            elif any(keyword in line_lower for keyword in ['**conclusion', 'conclusion', '**conclude', 'conclude']):
                current_section = "Conclusions"
                continue
            elif any(keyword in line_lower for keyword in ['**implication', 'implication', '**impact', 'impact', 'significance']):
                current_section = "Implications"
                continue
            elif current_section and line and len(line) > 10:
                # Clean up markdown formatting and bullet points
                clean_line = line.replace('**', '').replace('*', '').replace('•', '').replace('-', '').strip()
                if clean_line:
                    sections[current_section].append(clean_line)
        
        # If no clear sections were found, split content into paragraphs and distribute
        if all(len(section) == 0 for section in sections.values()):
            paragraphs = [p.strip() for p in summary.split('\n\n') if p.strip() and len(p.strip()) > 20]
            section_names = list(sections.keys())
            
            for i, paragraph in enumerate(paragraphs):
                section_index = i % len(section_names)
                sections[section_names[section_index]].append(paragraph)
        
        # Ensure each section has at least some content
        for section_name in sections:
            if not sections[section_name]:
                sections[section_name] = ["Content not explicitly separated in the original summary"]
                
        logging.info("Summary sections formatted successfully without additional AI processing")
        return sections
        
    except Exception as e:
        logging.error(f"Error in format_summary_sections: {str(e)}")
        # Fallback: put all content in Key Findings
        return {
            "Key Findings": [summary],
            "Methodology": ["See main summary for details"],
            "Conclusions": ["See main summary for details"], 
            "Implications": ["See main summary for details"]
        }

class SummaryPDF(FPDF):
    def header(self):
        self.set_font('Arial', 'B', 15)
        self.cell(0, 10, 'PaperSynth - Research Summary', 0, 1, 'C')
        self.ln(10)
        
    def footer(self):
        self.set_y(-15)
        self.set_font('Arial', 'I', 8)
        self.cell(0, 10, f'Generated on {datetime.now().strftime("%Y-%m-%d %H:%M")} - Page {self.page_no()}', 0, 0, 'C')

def save_summary_to_pdf(summary, output_path):
    try:
        pdf = SummaryPDF()
        pdf.add_page()
        
        # Format the summary into sections
        sections = format_summary_sections(summary)
        
        # Add formatted content using standard dash
        for section, points in sections.items():
            # Section header
            pdf.set_font('Arial', 'B', 14)
            pdf.cell(0, 10, section, ln=True)
            pdf.ln(5)
            
            # Section content
            pdf.set_font('Arial', '', 12)
            for point in points:
                # Clean the text of any problematic characters
                clean_point = point.encode('ascii', 'replace').decode()
                # Add content with dash instead of bullet
                pdf.multi_cell(0, 10, f"- {clean_point}")
            pdf.ln(5)
        
        pdf.output(output_path)
        logging.info(f"Summary saved to {output_path}")
        return output_path
    except Exception as e:
        logging.error(f"Error saving summary to PDF: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error saving summary to PDF: {str(e)}")
def generate_graphical_abstract(summary, pipe, request_dir, preset="balanced"):
    try:
        # Create a more focused prompt for research visualization
        prompt = f"""Create a scientific graphical abstract visualization:
        A clean, professional diagram showing:
        {summary[:300]}
        Style: Modern scientific illustration, minimalist, clear layout, professional colors
        Include: Relevant scientific symbols, data visualization elements, and clear visual hierarchy
        """

        negative_prompt = "text, words, blurry, low quality, distorted, messy, cluttered"

        # Choose parameters by preset
        if preset not in {"fast", "balanced", "quality"}:
            preset = "balanced"
        params = {
            "fast": {"steps": 20, "size": 384},
            "balanced": {"steps": 30, "size": 512},
            "quality": {"steps": 50, "size": 768},
        }[preset]

        try:
            output = pipe(
                prompt=prompt,
                negative_prompt=negative_prompt,
                num_inference_steps=params["steps"],
                guidance_scale=7.5,
                height=params["size"],
                width=params["size"],
                generator=torch.manual_seed(42)
            )
        except torch.cuda.OutOfMemoryError:
            # Fallback smaller/fewer steps
            fallback_size = 384 if params["size"] > 384 else 256
            fallback_steps = 20 if params["steps"] > 20 else 15
            output = pipe(
                prompt=prompt,
                negative_prompt=negative_prompt,
                num_inference_steps=fallback_steps,
                guidance_scale=7.5,
                height=fallback_size,
                width=fallback_size,
                generator=torch.manual_seed(42)
            )

        if not output.images:
            raise ValueError("No images generated")

        image = output.images[0]
        graphical_abstract_path = os.path.join(request_dir, "graphical_abstract.png")
        image.save(graphical_abstract_path, quality=95)
        return graphical_abstract_path
    except Exception as e:
        logging.error(f"Graphical Abstract Error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to generate graphical abstract: {str(e)}")

# AI Voiceover
def generate_voice(summary, request_dir):
    try:
        if not ElevenLabs:
            raise HTTPException(status_code=500, detail="ElevenLabs package not available. Install with: pip install elevenlabs")
        if not ELEVENLABS_API_KEY:
            raise HTTPException(status_code=500, detail="ElevenLabs API key not configured")
        
        client = ElevenLabs(api_key=ELEVENLABS_API_KEY)
        audio = client.generate(
            text=summary,
            voice="Lily",
            model="eleven_monolingual_v1"
        )
        voiceover_path = os.path.join(request_dir, "voiceover.mp3")
        with open(voiceover_path, "wb") as f:
            f.write(audio)
        return voiceover_path
    except Exception as e:
        logging.error(f"Voiceover Error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to generate voiceover: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Voiceover Error: {str(e)}")

# Generate Presentation
def generate_presentation(summary, request_dir):
    try:
        prs = Presentation()
        
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "PaperSynth - Research Summary"
        title_slide.placeholders[1].text = "Generated Summary Presentation"
        
        # Summary sections
        sections = format_summary_sections(summary)
        for section, points in sections.items():
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = section
            text_frame = slide.placeholders[1].text_frame
            
            for point in points:
                p = text_frame.add_paragraph()
                p.text = f"• {point}"
                p.level = 0
        
        presentation_path = os.path.join(request_dir, "presentation.pptx")
        prs.save(presentation_path)
        return presentation_path
    except Exception as e:
        logging.error(f"Presentation Generation Error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Presentation Generation Error: {str(e)}")

# FastAPI Route
@app.post("/process-paper/")
async def process_paper(
    request: Request,
    file: UploadFile,
    summary_length: str = "medium",
    generate_visual: bool = True,
    generate_audio: bool = True,
    sdxl_preset: str = "balanced",
    _: bool = Depends(require_bearer_token),
    __: bool = Depends(enforce_rate_limit),
):
    try:
        # Concurrency gate
        await _CONC_SEM.acquire()
        # Create per-request directory
        request_id = uuid.uuid4().hex
        request_dir = os.path.join(TEMP_DIR, request_id)
        os.makedirs(request_dir, exist_ok=True)
        
        # Initialize warnings list early
        warnings = []

        # Validate preset
        if sdxl_preset not in {"fast", "balanced", "quality"}:
            sdxl_preset = "balanced"

        # Determine which features to run (env gate + request flag)
        want_visual = ENABLE_SDXL and generate_visual
        want_audio = ENABLE_TTS and generate_audio

        # Lazy load Stable Diffusion only if needed
        pipe = None
        if want_visual:
            try:
                logging.info(f"[{request_id}] Starting SDXL lazy loading...")
                pipe = await lazy_load_stable_diffusion()
                logging.info(f"[{request_id}] SDXL loaded successfully")
            except Exception as e:
                logging.error(f"[{request_id}] Failed to load SDXL: {e}")
                want_visual = False  # Disable visual generation if loading fails
                warnings.append("SDXL_LOAD_FAILED: Could not load Stable Diffusion model")

        # Validate inputs
        if summary_length not in ALLOWED_SUMMARY_LENGTHS:
            raise HTTPException(status_code=400, detail=f"Invalid summary_length. Allowed: {sorted(ALLOWED_SUMMARY_LENGTHS)}")

        if not file or not file.filename or not file.filename.lower().endswith(".pdf"):
            raise HTTPException(status_code=400, detail="Please upload a .pdf file")

        if file.content_type and file.content_type not in ("application/pdf", "application/x-pdf", "application/acrobat", "applications/pdf", "text/pdf", "text/x-pdf"):
            raise HTTPException(status_code=400, detail="Uploaded file does not appear to be a PDF")

        file_path = os.path.join(request_dir, "paper.pdf")
        with open(file_path, "wb") as f:
            total = 0
            while True:
                chunk = file.file.read(1024 * 1024)
                if not chunk:
                    break
                total += len(chunk)
                if total > MAX_UPLOAD_BYTES:
                    raise HTTPException(status_code=413, detail="File too large. Max 10 MB")
                f.write(chunk)

        try:
            t0 = time.perf_counter()
            text, pdf_page_count = extract_text_from_pdf(file_path)
            logging.info(f"[{request_id}] pdf_extract_ms={(time.perf_counter()-t0)*1000:.0f} pages={pdf_page_count}")
        except HTTPException as he:
            code = "PDF_INVALID" if he.detail == "Invalid or corrupted PDF file" else "PDF_READ_FAILED"
            logging.error(f"[{request_id}] PDF error: {he.detail}")
            raise HTTPException(status_code=he.status_code, detail=_error_response(code, str(he.detail), request_id))

        try:
            t1 = time.perf_counter()
            summary = gemini_summary(text, summary_length)
            logging.info(f"[{request_id}] summarize_ms={(time.perf_counter()-t1)*1000:.0f}")
        except HTTPException as he:
            logging.error(f"[{request_id}] Summarization error: {he.detail}")
            raise HTTPException(status_code=he.status_code, detail=_error_response("SUMMARY_FAILED", "Failed to generate summary", request_id))
        
        # Generate all outputs
        summary_pdf_path = os.path.join(request_dir, "summary.pdf")
        t2 = time.perf_counter()
        save_summary_to_pdf(summary, summary_pdf_path)
        logging.info(f"[{request_id}] pdf_summary_ms={(time.perf_counter()-t2)*1000:.0f}")

        if want_visual and pipe is not None:
            try:
                t3 = time.perf_counter()
                graphical_abstract_path = generate_graphical_abstract(summary, pipe, request_dir, preset=sdxl_preset)
                logging.info(f"[{request_id}] sdxl_ms={(time.perf_counter()-t3)*1000:.0f}")
            except Exception as he:
                logging.error(f"[{request_id}] SDXL error: {getattr(he, 'detail', he)}")
                graphical_abstract_path = None
                warnings.append("SDXL_FAILED: Graphical abstract generation failed")
        else:
            graphical_abstract_path = None

        speaker_notes = None
        if want_audio:
            try:
                t4 = time.perf_counter()
                voiceover_path = generate_voice(summary, request_dir)
                logging.info(f"[{request_id}] tts_ms={(time.perf_counter()-t4)*1000:.0f}")
            except Exception as he:
                logging.error(f"[{request_id}] TTS error: {getattr(he, 'detail', he)}")
                voiceover_path = None
                speaker_notes = summary  # provide plain text as a fallback
                warnings.append("TTS_FAILED: Audio generation failed; provided speaker_notes instead")
        else:
            voiceover_path = None

        try:
            t5 = time.perf_counter()
            presentation_path = generate_presentation(summary, request_dir)
            logging.info(f"[{request_id}] pptx_ms={(time.perf_counter()-t5)*1000:.0f}")
        except HTTPException as he:
            logging.error(f"[{request_id}] PPTX error: {he.detail}")
            raise HTTPException(status_code=he.status_code, detail=_error_response("PRESENTATION_FAILED", "Failed to generate presentation", request_id))

        # Derive base URL from request (respects proxies if headers set)
        base_url = str(request.base_url).rstrip('/')
        def build_url(filename: Optional[str]):
            if not filename:
                return None
            if SIGNED_DOWNLOADS and DOWNLOAD_SIGNING_KEY:
                expires = int(time.time()) + 15 * 60  # 15 minutes
                payload = {"rid": request_id, "file": filename, "exp": expires}
                msg = f"{request_id}:{filename}:{expires}".encode()
                sig = hmac.new(DOWNLOAD_SIGNING_KEY.encode(), msg, hashlib.sha256).hexdigest()
                qs = urlencode({**payload, "sig": sig})
                return f"{base_url}/download?{qs}"
            else:
                return f"{base_url}/static/{request_id}/{filename}"

        return {
            "request_id": request_id,
            "summary": summary,
            "pages": pdf_page_count,
            "summary_pdf": build_url("summary.pdf"),
            "graphical_abstract": build_url("graphical_abstract.png"),
            "voiceover": build_url("voiceover.mp3"),
            "presentation": build_url("presentation.pptx"),
            "features": {"sdxl": ENABLE_SDXL, "tts": ENABLE_TTS, "signed_downloads": SIGNED_DOWNLOADS},
            "speaker_notes": speaker_notes,
            "warnings": warnings,
        }
    except Exception as e:
        logging.error(f"[{locals().get('request_id','-')}] Error processing paper: {str(e)}")
        # If it's already an HTTPException, pass through
        if isinstance(e, HTTPException):
            raise e
        raise HTTPException(status_code=500, detail=_error_response("INTERNAL_ERROR", "Unexpected error", locals().get("request_id")))
    finally:
        # Resource cleanup and monitoring
        try:
            request_id = locals().get('request_id', 'unknown')
            
            # Log memory usage for monitoring
            memory_usage = get_memory_usage()
            logging.info(f"[{request_id}] Request completed. Memory: {memory_usage['rss_mb']}MB ({memory_usage['percent']}%)")
            
            # Cleanup GPU memory if we used SDXL
            if locals().get('want_visual') and locals().get('pipe'):
                cleanup_gpu_memory()
            
            # Release concurrency slot
            _CONC_SEM.release()
        except Exception as cleanup_error:
            logging.warning(f"Cleanup error: {cleanup_error}")

# Run FastAPI Server
if __name__ == "__main__":
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
